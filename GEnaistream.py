import streamlit as st
import pandas as pd
import PyPDF2
import textract
import os
import pptx
import GenAi

def read_pdf(file):
    text = ""
    with open(file, "rb") as f:
        pdf_reader = PyPDF2.PdfFileReader(f)
        for page_num in range(pdf_reader.numPages):
            page = pdf_reader.getPage(page_num)
            text += page.extractText()
    return text

def read_docx(file):
    text = ""
    docx_text = textract.process(file)
    text += docx_text.decode("utf-8")
    return text

def read_ppt(file):
    ppt_text = ""
    prs = pptx.Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += shape.text
    return ppt_text

def main():
    st.title("File Summarizer")

    uploaded_file = st.file_uploader("Upload a file ", type=["csv", "txt", "pdf", "pptx", "docx"])

    if uploaded_file is not None:
        st.write("File uploaded successfully")

        file_extension = os.path.splitext(uploaded_file.name)[1].lower()

        if file_extension == ".csv":
            df = pd.read_csv(uploaded_file)
            option=st.selectbox("Choose your style",options=("casual","Formal"))
            GenAi.getsummary(df,option)
            st.write(df,option)
        elif file_extension == ".txt":
            text = uploaded_file.getvalue().decode("utf-8")
            GenAi.getsummary(text)
            st.write(text)
        elif file_extension == ".pdf":
            pdf_text = read_pdf(uploaded_file)
            GenAi.getsummary(pdf_text)
            st.write(pdf_text)
        elif file_extension == ".pptx":
            ppt_text = read_ppt(uploaded_file)
            GenAi.getsummary(ppt_text)
            st.write(ppt_text)
        elif file_extension == ".docx":
            docx_text = read_docx(uploaded_file)
            st.write(docx_text)
        else:
            st.write("Unsupported file type")

if __name__ == "__main__":
    main()